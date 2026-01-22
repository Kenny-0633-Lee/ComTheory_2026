import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# 파일명에서 챕터 번호만 추출하는 정규식 (예: Chapter02a -> 02)
# 대소문자 무시, Chapter 뒤의 숫자 그룹을 캡처
FILENAME_PATTERN = re.compile(r"chapter(\d+)[a-z]*", re.IGNORECASE)

def get_chapter_number(filename):
    match = FILENAME_PATTERN.search(filename)
    if match:
        return match.group(1) # '02', '11' 등 문자열 반환
    return "unknown"

def process_chapter_group(ch_num, file_list, source_dir, base_output_dir):
    # 출력 폴더 생성 (예: assets_extracted/ch02)
    ch_dir_name = f"ch{ch_num}"
    ch_output_dir = os.path.join(base_output_dir, ch_dir_name)
    os.makedirs(ch_output_dir, exist_ok=True)
    
    log_file = os.path.join(ch_output_dir, f"captions_{ch_dir_name}.txt")
    
    # 이미지 번호 카운터 (파일이 바뀌어도 유지됨)
    img_global_idx = 1
    
    print(f"📂 Processing Chapter {ch_num} ({len(file_list)} files) -> {ch_output_dir}")

    with open(log_file, "w", encoding="utf-8") as log:
        # 파일 목록을 알파벳 순으로 정렬 (Chapter02a -> Chapter02b)
        for pptx_file in sorted(file_list):
            pptx_path = os.path.join(source_dir, pptx_file)
            print(f"  - Reading {pptx_file}...")
            
            try:
                prs = Presentation(pptx_path)
            except Exception as e:
                print(f"    ❌ Error reading {pptx_file}: {e}")
                continue

            for i, slide in enumerate(prs.slides):
                # 캡션 찾기
                caption_text = "No Caption Found"
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        t = shape.text_frame.text.strip()
                        # "Figure" 혹은 "Fig"로 시작하는 텍스트 감지
                        if t.lower().startswith("figure") or t.lower().startswith("fig"):
                            caption_text = t
                            break # 캡션 하나 찾으면 중단 (보통 슬라이드당 하나)
                
                # 이미지 추출
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        # 파일명: fig_ch02_01.png (파일이 a/b로 나뉘어도 번호는 계속 증가)
                        img_filename = f"fig_{ch_dir_name}_{img_global_idx:02d}.{image.ext}"
                        img_path = os.path.join(ch_output_dir, img_filename)
                        
                        with open(img_path, "wb") as f:
                            f.write(image.blob)
                        
                        # 로그에 원본 파일명도 같이 기록해두면 나중에 추적하기 좋습니다.
                        log.write(f"[{img_filename}] Source: {pptx_file} | Slide: {i+1} | Caption: {caption_text}\n")
                        img_global_idx += 1

if __name__ == "__main__":
    SOURCE_DIR = "source_pptx"
    EXTRACT_DIR = "assets_extracted"
    
    if not os.path.exists(SOURCE_DIR):
        print(f"Error: '{SOURCE_DIR}' directory not found.")
        exit()

    # 1. 파일 스캔 및 그룹화
    # groups = { '02': ['Chapter02a.pptx', 'Chapter02b.pptx'], '03': ... }
    groups = {}
    
    all_files = os.listdir(SOURCE_DIR)
    
    # .ppt 파일 경고
    ppt_files = [f for f in all_files if f.lower().endswith(".ppt")]
    if ppt_files:
        print("\n⚠️  WARNING: Found .ppt files. These must be converted to .pptx!")
        print(f"   Files: {ppt_files[:3]} ... and others.\n")

    # .pptx 파일만 처리
    pptx_files = [f for f in all_files if f.lower().endswith(".pptx")]
    
    for f in pptx_files:
        ch_num = get_chapter_number(f)
        if ch_num not in groups:
            groups[ch_num] = []
        groups[ch_num].append(f)

    # 2. 그룹별 처리 실행
    # 챕터 번호 순서대로 실행 (01, 02, 03 ...)
    for ch_num in sorted(groups.keys()):
        process_chapter_group(ch_num, groups[ch_num], SOURCE_DIR, EXTRACT_DIR)
        
    print("\n✅ All extraction tasks completed.")